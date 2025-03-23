#========================= Imports ===========================#

import os
import re
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from arabic_reshaper import reshape
from peft import PeftModel
import torch
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx import Presentation
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer



#======================== Load model and tokenizer ====================#

base_model = AutoModelForSeq2SeqLM.from_pretrained("Helsinki-NLP/opus-mt-tc-big-en-ar")
tokenizer = AutoTokenizer.from_pretrained("./finetuned_lora_model")
lora_model = PeftModel.from_pretrained(base_model, "./finetuned_lora_model")




#======================= Function to translate text =========================#

def translate_text(text):
    if re.match(r'^[0-9IVXLCDM]+$', text.strip()):
        return text
    if not text.strip():
        return text
    inputs = tokenizer(text, return_tensors="pt", padding=True, truncation=True, max_length=128)
    with torch.no_grad():
        outputs = lora_model.generate(**inputs, max_length=128)
    arabic_text = tokenizer.batch_decode(outputs, skip_special_tokens=True)[0]
    arabic_text = reshape(arabic_text)
    return arabic_text



#======================= Function to rotate arrow shapes by 180 degree =====================#

def rotate_shape(slide, shape):
    if shape.auto_shape_type in [MSO_SHAPE.RECTANGLE, MSO_SHAPE.ISOSCELES_TRIANGLE, MSO_SHAPE.ROUNDED_RECTANGLE]:
        return
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text  
        font = shape.text_frame.paragraphs[0].font  
        font_size = font.size or Pt(16)  
        alignment = shape.text_frame.paragraphs[0].alignment 
        if alignment == PP_ALIGN.LEFT:
            alignment = PP_ALIGN.RIGHT 

        try:
            font_color = font.color.rgb  
        except AttributeError:
            font_color = RGBColor(255, 255, 255)  

        shape.text_frame.clear()
        shape.rotation = (shape.rotation + 180) % 360

        new_textbox = slide.shapes.add_textbox(left, top, width, height)
        new_textbox.text_frame.paragraphs[0].font.language_id  = MSO_LANGUAGE_ID.ARABIC
        new_textbox.text_frame.text = text
        new_textbox.text_frame.paragraphs[0].font.size = font_size
        new_textbox.text_frame.paragraphs[0].alignment = alignment
        new_textbox.text_frame.paragraphs[0].font.color.rgb = font_color 
        new_textbox.text_frame.word_wrap = True

    else:
        shape.rotation = (shape.rotation + 180) % 360



#==================== Function to swap columns of a table =======================#

def reverse_table_columns(table):
    num_cols = len(table.columns)
    num_rows = len(table.rows)

    has_merged_cells = False
    for row in range(num_rows):
        for col in range(num_cols):
            cell = table.cell(row, col)
            if cell._tc is not None:
                grid_span = cell._tc.get("gridSpan")
                if grid_span:
                    has_merged_cells = True
                    break
        if has_merged_cells:
            break

    if has_merged_cells:
        for row in range(num_rows):
            for col in range(num_cols):
                cell = table.cell(row, col)
                cell_text = cell.text
                cell.text = cell_text
                for paragraph in cell.text_frame.paragraphs:
                    text = translate_text(paragraph.text)
                    paragraph.text = text
            
                    paragraph.font.language_id = MSO_LANGUAGE_ID.ARABIC
                    paragraph.font.bold = paragraph.font.bold
                    paragraph.font.italic = paragraph.font.italic
                    paragraph.font.underline = paragraph.font.underline
                    if paragraph.alignment in [PP_ALIGN.LEFT, None]:
                        paragraph.alignment = PP_ALIGN.RIGHT

        return

    for row in range(num_rows):
        values = [table.cell(row, col).text for col in range(num_cols)]
        for col in range(num_cols):
            cell_text = values[num_cols - 1 - col]  
            cell = table.cell(row, col)
            cell.text = cell_text

            for paragraph in cell.text_frame.paragraphs:
                text = translate_text(paragraph.text)
                paragraph.text = text
            
                paragraph.font.language_id = MSO_LANGUAGE_ID.ARABIC
                paragraph.font.bold = paragraph.font.bold
                paragraph.font.italic = paragraph.font.italic
                paragraph.font.underline = paragraph.font.underline
                if paragraph.alignment in [PP_ALIGN.LEFT, None]:
                    paragraph.alignment = PP_ALIGN.RIGHT

    column_widths = [table.columns[col].width for col in range(num_cols)]
    for col in range(num_cols):
        table.columns[col].width = column_widths[num_cols - 1 - col]




# ================================ Function to mirror each shape of a slide ======================= #

def mirror_shape(slide, shape, slide_width):
    if shape.left is not None:
        if shape.has_table:
            original_left = shape.left
            shape_width = shape.width
            new_left = slide_width - (original_left + shape_width)
            shape.left = new_left
            shape.width = shape_width
            reverse_table_columns(shape.table)
        
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            rotate_shape(slide, shape)
            original_left = shape.left
            shape_width = shape.width
            new_left = slide_width - (original_left + shape_width)
            shape.left = new_left
            shape.width = shape_width
            
        else:
            original_left = shape.left
            shape_width = shape.width
            new_left = slide_width - (original_left + shape_width)
            shape.left = new_left
            shape.width = shape_width

    if shape.has_text_frame and shape.text_frame.text.strip():  
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text
            translated_text = translate_text(text)
            paragraph.text = translated_text
            paragraph.font.language_id = MSO_LANGUAGE_ID.ARABIC
            paragraph.font.bold = paragraph.font.bold
            paragraph.font.italic = paragraph.font.italic
            paragraph.font.underline = paragraph.font.underline
            if paragraph.alignment in [PP_ALIGN.LEFT, None]:
                paragraph.alignment = PP_ALIGN.RIGHT

        shape.text_frame.word_wrap = True


#===================== Function to find and mirror grouped shape ========================#

def mirror_grouped_shape(slide, group, slide_width):
    group_left = min([shape.left for shape in group.shapes])
    new_group_left = slide_width - group_left
    shift_x = new_group_left - group_left
    for shape in group.shapes:
        shape.left += shift_x
        mirror_shape(slide, shape, slide_width)
    group.left = new_group_left



#=================== Function to mirror the whole slide ========================#

def mirror_slide(slide, slide_width, progress_callback, current_slide, total_slides):
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.has_text_frame and shape.text_frame.text.strip():
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text
                    translated_text = translate_text(text)
                    paragraph.text = translated_text
                    paragraph.font.language_id = MSO_LANGUAGE_ID.ARABIC
                    paragraph.font.bold = paragraph.font.bold
                    paragraph.font.italic = paragraph.font.italic
                    paragraph.font.underline = paragraph.font.underline
                    if paragraph.alignment in [PP_ALIGN.LEFT, None]:
                        paragraph.alignment = PP_ALIGN.RIGHT

                shape.text_frame.word_wrap = True
            continue

        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            mirror_grouped_shape(slide, shape, slide_width)

        else:
            mirror_shape(slide, shape, slide_width)

    progress_callback(current_slide, total_slides)
    




#========================= Main Function ========================#
def modifiy_presentation(input_pptx, output_pptx, progress_callback):
    prs = Presentation(input_pptx)
    slide_width = prs.slide_width
    total_slides = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        mirror_slide(slide, slide_width, progress_callback, idx+1, total_slides)

    prs.save(output_pptx)



#========================= Streamlit UI =======================#
os.makedirs("app-processing-files", exist_ok=True)

st.title("PowerPoint Translator (English to Arabic)")

uploaded_file = st.file_uploader("Upload a PowerPoint file (.pptx)", type=["pptx"])

if uploaded_file is not None:
    with open("app-processing-files/uploaded_presentation.pptx", "wb") as f:
        f.write(uploaded_file.getvalue())

    progress_bar = st.progress(0)
    progress_text = st.empty()

    def update_progress(current, total):
        percentage = int((current / total) * 100)
        progress_bar.progress(percentage / 100)
        progress_text.text(f"Processing Slide {current}/{total} ({percentage}%)")

    output_file = "app-processing-files/translated_presentation.pptx"

    if st.button("Translate Presentation"):
        modifiy_presentation("app-processing-files/uploaded_presentation.pptx", output_file, update_progress)

        with open(output_file, "rb") as f:
            st.download_button(
                label="Download Translated Presentation",
                data=f,
                file_name="translated_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        st.success("Translation complete! You can download the file above.")